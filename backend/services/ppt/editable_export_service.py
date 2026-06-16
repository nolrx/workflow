"""
Editable PPTX Export Service - creates PPTX with editable text boxes

This service implements the complete flow:
1. OCR layout analysis - detect text elements and bounding boxes
2. Style extraction - use AI to extract text attributes (color, bold, italic, etc.)
3. Background removal - inpainting to remove text from original image
4. PPTX reconstruction - create slides with clean background + editable text boxes
"""
import json
import logging
import os
from dataclasses import dataclass, field
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches

from backend.services.ppt.prompts import (
    get_batch_text_attribute_extraction_prompt,
    get_clean_background_prompt,
    get_quality_enhancement_prompt,
    get_text_attribute_extraction_prompt,
)

logger = logging.getLogger(__name__)

# Standard 16:9 slide dimensions
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Image resolution (4K = 3840x2160, but we use 1920x1080 for performance)
IMAGE_WIDTH = 1920
IMAGE_HEIGHT = 1080


@dataclass
class TextElement:
    """Represents a detected text element with position and style"""
    id: str
    text: str
    bbox: Tuple[int, int, int, int]  # (x, y, width, height) in pixels
    # Style attributes (filled by AI)
    font_color: str = "#000000"
    is_bold: bool = False
    is_italic: bool = False
    is_underline: bool = False
    text_alignment: str = "left"  # left, center, right, justify
    # Optional colored segments for multi-color text
    colored_segments: List[Dict] = field(default_factory=list)


@dataclass
class ExportResult:
    """Result of editable export operation"""
    success: bool
    pptx_path: Optional[str] = None
    error: Optional[str] = None
    pages_processed: int = 0
    text_elements_found: int = 0


def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return (0, 0, 0)  # Default to black
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except ValueError:
        return (0, 0, 0)


def get_alignment_enum(alignment: str):
    """Convert alignment string to python-pptx enum"""
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY,
    }
    return alignment_map.get(alignment.lower(), PP_ALIGN.LEFT)


class EditableExportService:
    """
    Service for creating editable PPTX exports.

    Workflow:
    1. For each page image:
       a. Run OCR to detect text elements and positions
       b. Use AI to extract text attributes (colors, bold, etc.)
       c. Generate clean background (remove text via inpainting)
       d. Create slide with background image + positioned text boxes
    """

    def __init__(
        self,
        ai_provider,
        ocr_provider=None,
        inpaint_method: str = "generative",
        extractor_method: str = "hybrid"
    ):
        """
        Initialize the export service.

        Args:
            ai_provider: AI provider for image generation/analysis
            ocr_provider: OCR provider for text detection (optional, uses AI if None)
            inpaint_method: Background removal method - 'baidu', 'generative', or 'hybrid'
            extractor_method: Text extraction method - 'ocr', 'ai', or 'hybrid'
        """
        self.ai_provider = ai_provider
        self.ocr_provider = ocr_provider
        self.inpaint_method = inpaint_method
        self.extractor_method = extractor_method

    def export_project(
        self,
        project_id: str,
        page_images: List[Tuple[str, bytes]],
        output_path: str,
        page_descriptions: Optional[List[str]] = None
    ) -> ExportResult:
        """
        Export project pages to editable PPTX.

        Args:
            project_id: Project ID for logging
            page_images: List of (page_id, image_bytes) tuples
            output_path: Path to save the PPTX file
            page_descriptions: Optional list of page descriptions

        Returns:
            ExportResult with success status and details
        """
        if not page_images:
            return ExportResult(success=False, error="No pages to export")

        try:
            prs = Presentation()
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = SLIDE_HEIGHT
            blank_layout = prs.slide_layouts[6]

            total_text_elements = 0

            for i, (page_id, image_data) in enumerate(page_images):
                logger.info(f"Processing page {i + 1}/{len(page_images)}: {page_id}")

                # Get page description if available
                description = page_descriptions[i] if page_descriptions and i < len(page_descriptions) else None

                # Process single page
                slide, text_count = self._process_page(
                    prs=prs,
                    layout=blank_layout,
                    image_data=image_data,
                    page_index=i,
                    description=description
                )

                total_text_elements += text_count

            # Ensure output directory exists
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            # Save PPTX
            prs.save(output_path)

            return ExportResult(
                success=True,
                pptx_path=output_path,
                pages_processed=len(page_images),
                text_elements_found=total_text_elements
            )

        except Exception as e:
            logger.error(f"Editable export failed: {str(e)}", exc_info=True)
            return ExportResult(success=False, error=str(e))

    def _process_page(
        self,
        prs: Presentation,
        layout,
        image_data: bytes,
        page_index: int,
        description: Optional[str] = None
    ) -> Tuple[Any, int]:
        """
        Process a single page for editable export.

        Returns:
            Tuple of (slide, text_elements_count)
        """
        # Step 1: Detect text elements
        text_elements = self._detect_text_elements(image_data, page_index)

        # Step 2: Extract text attributes using AI
        if text_elements:
            text_elements = self._extract_text_attributes(image_data, text_elements)

        # Step 3: Generate clean background
        clean_background = self._generate_clean_background(
            image_data,
            text_elements
        )

        # Step 4: Create slide with background and text boxes
        slide = self._create_slide(
            prs=prs,
            layout=layout,
            background_image=clean_background or image_data,
            text_elements=text_elements
        )

        return slide, len(text_elements)

    def _detect_text_elements(
        self,
        image_data: bytes,
        page_index: int
    ) -> List[TextElement]:
        """
        Detect text elements and their positions in the image.

        Uses OCR provider if available, otherwise uses AI vision.
        """
        if self.ocr_provider:
            return self._detect_with_ocr(image_data)
        else:
            return self._detect_with_ai(image_data, page_index)

    def _detect_with_ocr(self, image_data: bytes) -> List[TextElement]:
        """Use OCR provider for text detection"""
        # This would integrate with MinerU, Baidu OCR, or similar
        # For now, return empty list - to be implemented with actual OCR integration
        logger.warning("OCR provider not fully implemented, returning empty elements")
        return []

    def _detect_with_ai(self, image_data: bytes, page_index: int) -> List[TextElement]:
        """
        Use AI vision to detect text elements.

        This is a fallback when OCR is not available.
        """
        prompt = """
Analyze this PPT slide image and identify all text elements.
For each text element, provide:
1. The exact text content
2. Approximate position (x, y, width, height) as percentages of image dimensions

Return as JSON array:
[
    {
        "id": "0",
        "text": "Slide Title",
        "bbox": {"x_pct": 10, "y_pct": 5, "width_pct": 80, "height_pct": 15}
    },
    ...
]

Only output JSON, no other text.
"""
        try:
            result = self.ai_provider.generate_text(
                prompt,
                reference_images=[image_data]
            )

            if not result.success:
                logger.warning(f"AI text detection failed: {result.error}")
                return []

            # Parse response
            response_text = result.text.strip()
            if response_text.startswith('```'):
                response_text = response_text.strip('`').strip()
                if response_text.startswith('json'):
                    response_text = response_text[4:].strip()

            data = json.loads(response_text)

            # Convert to TextElement objects
            elements = []
            for item in data:
                bbox = item.get('bbox', {})
                # Convert percentages to pixels (assuming 1920x1080)
                x = int(bbox.get('x_pct', 0) * IMAGE_WIDTH / 100)
                y = int(bbox.get('y_pct', 0) * IMAGE_HEIGHT / 100)
                w = int(bbox.get('width_pct', 10) * IMAGE_WIDTH / 100)
                h = int(bbox.get('height_pct', 5) * IMAGE_HEIGHT / 100)

                elements.append(TextElement(
                    id=str(item.get('id', len(elements))),
                    text=item.get('text', ''),
                    bbox=(x, y, w, h)
                ))

            return elements

        except Exception as e:
            logger.error(f"AI text detection error: {str(e)}")
            return []

    def _extract_text_attributes(
        self,
        image_data: bytes,
        text_elements: List[TextElement]
    ) -> List[TextElement]:
        """
        Extract text attributes (color, bold, etc.) using AI.

        Uses hybrid strategy:
        - Full image analysis for layout attributes (bold, italic, alignment)
        - Cropped regions for precise color extraction
        """
        if not text_elements:
            return text_elements

        # Step 1: Batch extraction for layout attributes
        elements_for_prompt = [
            {"id": elem.id, "text": elem.text}
            for elem in text_elements
        ]

        batch_prompt = get_batch_text_attribute_extraction_prompt(elements_for_prompt)

        try:
            result = self.ai_provider.generate_text(
                batch_prompt,
                reference_images=[image_data]
            )

            if result.success:
                # Parse batch attributes
                response_text = result.text.strip()
                if response_text.startswith('```'):
                    response_text = response_text.strip('`').strip()
                    if response_text.startswith('json'):
                        response_text = response_text[4:].strip()

                batch_attrs = json.loads(response_text)

                # Apply batch attributes to elements
                for elem in text_elements:
                    attrs = batch_attrs.get(elem.id, {})
                    elem.is_bold = attrs.get('is_bold', False)
                    elem.is_italic = attrs.get('is_italic', False)
                    elem.is_underline = attrs.get('is_underline', False)
                    elem.text_alignment = attrs.get('text_alignment', 'left')
                    elem.font_color = attrs.get('font_color', '#000000')

        except Exception as e:
            logger.warning(f"Batch attribute extraction failed: {str(e)}")

        # Step 2: Individual color extraction for hybrid mode
        if self.extractor_method == 'hybrid':
            text_elements = self._extract_precise_colors(image_data, text_elements)

        return text_elements

    def _extract_precise_colors(
        self,
        image_data: bytes,
        text_elements: List[TextElement]
    ) -> List[TextElement]:
        """
        Extract precise colors by cropping individual text regions.

        Full-image color detection is often inaccurate, so we crop each
        text region and analyze it separately.
        """
        try:
            image = Image.open(BytesIO(image_data))

            for elem in text_elements:
                x, y, w, h = elem.bbox

                # Add padding
                padding = 5
                x1 = max(0, x - padding)
                y1 = max(0, y - padding)
                x2 = min(image.width, x + w + padding)
                y2 = min(image.height, y + h + padding)

                # Crop region
                cropped = image.crop((x1, y1, x2, y2))

                # Convert to bytes
                buffer = BytesIO()
                cropped.save(buffer, format='PNG')
                cropped_bytes = buffer.getvalue()

                # Get color extraction prompt
                prompt = get_text_attribute_extraction_prompt(elem.text)

                try:
                    result = self.ai_provider.generate_text(
                        prompt,
                        reference_images=[cropped_bytes]
                    )

                    if result.success:
                        response_text = result.text.strip()
                        if response_text.startswith('```'):
                            response_text = response_text.strip('`').strip()
                            if response_text.startswith('json'):
                                response_text = response_text[4:].strip()

                        color_data = json.loads(response_text)
                        segments = color_data.get('colored_segments', [])

                        if segments:
                            elem.colored_segments = segments
                            # Use first segment color as primary
                            elem.font_color = segments[0].get('color', elem.font_color)

                except Exception as e:
                    logger.debug(f"Color extraction failed for element {elem.id}: {str(e)}")

        except Exception as e:
            logger.warning(f"Precise color extraction failed: {str(e)}")

        return text_elements

    def _generate_clean_background(
        self,
        image_data: bytes,
        text_elements: List[TextElement]
    ) -> Optional[bytes]:
        """
        Generate clean background by removing text.

        Uses the configured inpaint method:
        - 'baidu': Mask-based removal using Baidu API
        - 'generative': AI-powered regeneration
        - 'hybrid': Baidu removal + AI quality enhancement
        """
        if self.inpaint_method == 'baidu':
            return self._inpaint_with_baidu(image_data, text_elements)
        elif self.inpaint_method == 'generative':
            return self._inpaint_with_ai(image_data)
        else:  # hybrid
            result = self._inpaint_with_baidu(image_data, text_elements)
            if result:
                return self._enhance_quality(result, text_elements)
            return self._inpaint_with_ai(image_data)

    def _inpaint_with_baidu(
        self,
        image_data: bytes,
        text_elements: List[TextElement]
    ) -> Optional[bytes]:
        """
        Remove text using Baidu inpainting API.

        Creates a mask of text regions and uses API to fill them.
        """
        # This would integrate with Baidu's image inpainting API
        # For now, fall back to AI-based approach
        logger.info("Baidu inpainting not implemented, falling back to AI")
        return None

    def _inpaint_with_ai(self, image_data: bytes) -> Optional[bytes]:
        """
        Remove text using AI generative model.
        """
        prompt = get_clean_background_prompt()

        try:
            result = self.ai_provider.generate_image(
                prompt=prompt,
                reference_images=[image_data]
            )

            if result.success and result.image_data:
                return result.image_data
            else:
                logger.warning(f"AI inpainting failed: {result.error}")
                return None

        except Exception as e:
            logger.error(f"AI inpainting error: {str(e)}")
            return None

    def _enhance_quality(
        self,
        image_data: bytes,
        text_elements: List[TextElement]
    ) -> bytes:
        """
        Enhance quality of inpainted regions to remove artifacts.
        """
        # Build region list for the prompt
        regions = [
            {
                'x': elem.bbox[0],
                'y': elem.bbox[1],
                'width': elem.bbox[2],
                'height': elem.bbox[3]
            }
            for elem in text_elements
        ]

        prompt = get_quality_enhancement_prompt(regions)

        try:
            result = self.ai_provider.generate_image(
                prompt=prompt,
                reference_images=[image_data]
            )

            if result.success and result.image_data:
                return result.image_data
            else:
                logger.warning(f"Quality enhancement failed: {result.error}")
                return image_data

        except Exception as e:
            logger.warning(f"Quality enhancement error: {str(e)}")
            return image_data

    def _create_slide(
        self,
        prs: Presentation,
        layout,
        background_image: bytes,
        text_elements: List[TextElement]
    ) -> Any:
        """
        Create a slide with background image and editable text boxes.
        """
        slide = prs.slides.add_slide(layout)

        # Add background image
        image_stream = BytesIO(background_image)
        slide.shapes.add_picture(
            image_stream,
            left=0,
            top=0,
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT
        )

        # Add text boxes for each element
        for elem in text_elements:
            self._add_text_box(slide, elem)

        return slide

    def _add_text_box(self, slide, elem: TextElement):
        """
        Add an editable text box to the slide.
        """
        # Convert pixel coordinates to EMUs
        # Assuming image is 1920x1080, slide is 10" x 5.625"
        x_ratio = SLIDE_WIDTH / IMAGE_WIDTH
        y_ratio = SLIDE_HEIGHT / IMAGE_HEIGHT

        left = Emu(int(elem.bbox[0] * x_ratio))
        top = Emu(int(elem.bbox[1] * y_ratio))
        width = Emu(int(elem.bbox[2] * x_ratio))
        height = Emu(int(elem.bbox[3] * y_ratio))

        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        # Handle multi-color text (colored segments)
        if elem.colored_segments:
            p = tf.paragraphs[0]
            p.alignment = get_alignment_enum(elem.text_alignment)

            for segment in elem.colored_segments:
                run = p.add_run()
                run.text = segment.get('text', '')

                # Set font properties
                font = run.font
                color = segment.get('color', elem.font_color)
                r, g, b = hex_to_rgb(color)
                font.color.rgb = RGBColor(r, g, b)
                font.bold = elem.is_bold
                font.italic = elem.is_italic
                font.underline = elem.is_underline

        else:
            # Simple single-color text
            p = tf.paragraphs[0]
            p.text = elem.text
            p.alignment = get_alignment_enum(elem.text_alignment)

            # Set font properties
            font = p.runs[0].font if p.runs else p.font
            r, g, b = hex_to_rgb(elem.font_color)
            font.color.rgb = RGBColor(r, g, b)
            font.bold = elem.is_bold
            font.italic = elem.is_italic
            font.underline = elem.is_underline


def create_editable_pptx(
    ai_provider,
    project_id: str,
    page_images: List[Tuple[str, bytes]],
    output_path: str,
    inpaint_method: str = "generative",
    extractor_method: str = "hybrid",
    page_descriptions: Optional[List[str]] = None
) -> ExportResult:
    """
    Convenience function to create editable PPTX.

    Args:
        ai_provider: AI provider instance
        project_id: Project ID
        page_images: List of (page_id, image_bytes) tuples
        output_path: Output file path
        inpaint_method: Background removal method
        extractor_method: Text extraction method
        page_descriptions: Optional page descriptions

    Returns:
        ExportResult
    """
    service = EditableExportService(
        ai_provider=ai_provider,
        inpaint_method=inpaint_method,
        extractor_method=extractor_method
    )

    return service.export_project(
        project_id=project_id,
        page_images=page_images,
        output_path=output_path,
        page_descriptions=page_descriptions
    )
