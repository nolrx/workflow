"""
PPT Export Service - handles PPTX and PDF export
"""
import os
import logging
from pathlib import Path
from typing import List, Optional
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

# Standard 16:9 slide dimensions
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


def create_simple_pptx(
    image_paths: List[str],
    output_path: str,
    slide_width: float = SLIDE_WIDTH,
    slide_height: float = SLIDE_HEIGHT
) -> str:
    """
    Create a simple PPTX with images as full-slide backgrounds.

    Args:
        image_paths: List of paths to generated page images
        output_path: Path to save the PPTX file
        slide_width: Slide width in EMUs (default 16:9)
        slide_height: Slide height in EMUs (default 16:9)

    Returns:
        Path to the created PPTX file
    """
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Use blank layout (index 6)
    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:
        if not os.path.exists(image_path):
            logger.warning(f"Image not found, skipping: {image_path}")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add image as full slide background
        slide.shapes.add_picture(
            image_path,
            left=0,
            top=0,
            width=slide_width,
            height=slide_height
        )

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    prs.save(output_path)
    logger.info(f"Created PPTX with {len(image_paths)} slides: {output_path}")

    return output_path


def create_pptx_from_bytes(
    image_data_list: List[bytes],
    output_path: str,
    slide_width: float = SLIDE_WIDTH,
    slide_height: float = SLIDE_HEIGHT
) -> str:
    """
    Create a PPTX from image bytes (for in-memory processing).

    Args:
        image_data_list: List of image bytes
        output_path: Path to save the PPTX file
        slide_width: Slide width
        slide_height: Slide height

    Returns:
        Path to the created PPTX file
    """
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    blank_layout = prs.slide_layouts[6]

    for image_data in image_data_list:
        if not image_data:
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Create BytesIO from image data
        image_stream = BytesIO(image_data)

        slide.shapes.add_picture(
            image_stream,
            left=0,
            top=0,
            width=slide_width,
            height=slide_height
        )

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    return output_path


def export_project_to_pptx(
    project_id: str,
    pages: List,
    upload_folder: Path,
    filename: Optional[str] = None
) -> str:
    """
    Export a PPT project to PPTX format.

    Args:
        project_id: Project ID
        pages: List of PPTPage objects or dicts
        upload_folder: Base upload folder path
        filename: Optional custom filename

    Returns:
        Path to the exported PPTX file (relative to upload_folder)
    """
    # Collect image paths from pages
    image_paths = []
    pages_dir = upload_folder / "ppt" / project_id / "pages"

    for page in pages:
        # Handle both dict and object access
        if hasattr(page, 'generated_image_path'):
            image_path = page.generated_image_path
        else:
            image_path = page.get('generated_image_path')

        if image_path:
            # Convert relative path to absolute
            if image_path.startswith('/'):
                # Remove leading slash for path joining
                image_path = image_path.lstrip('/')

            # Try multiple path formats
            full_path = None

            # Format 1: Direct path under upload folder
            candidate = upload_folder / image_path.lstrip('/')
            if candidate.exists():
                full_path = str(candidate)

            # Format 2: Path is just filename, look in pages dir
            if not full_path:
                candidate = pages_dir / os.path.basename(image_path)
                if candidate.exists():
                    full_path = str(candidate)

            # Format 3: Full path from generated_image_path
            if not full_path and os.path.exists(image_path):
                full_path = image_path

            if full_path:
                image_paths.append(full_path)
            else:
                logger.warning(f"Could not find image for page: {image_path}")

    if not image_paths:
        raise ValueError("No images found to export")

    # Generate output path
    exports_dir = upload_folder / "ppt" / project_id / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)

    if not filename:
        filename = f"presentation_{project_id[:8]}.pptx"
    elif not filename.endswith('.pptx'):
        filename += '.pptx'

    output_path = str(exports_dir / filename)

    # Create PPTX
    create_simple_pptx(image_paths, output_path)

    # Return relative path for storage
    return f"/{project_id}/exports/{filename}"


def get_export_file_path(
    project_id: str,
    filename: str,
    upload_folder: Path
) -> Optional[str]:
    """
    Get the full path to an exported file.

    Args:
        project_id: Project ID
        filename: Export filename
        upload_folder: Base upload folder

    Returns:
        Full path if exists, None otherwise
    """
    exports_dir = upload_folder / "ppt" / project_id / "exports"
    filepath = exports_dir / filename

    if filepath.exists():
        return str(filepath)
    return None
