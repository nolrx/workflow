"""
PPT Export routes - handles PPTX and PDF export
"""
import logging
import subprocess
import tempfile
from pathlib import Path

from flask import Blueprint, current_app, request, send_file
from flask_jwt_extended import get_jwt_identity, jwt_required

from backend.models.ppt import PPTPage, PPTProject
from backend.services.ppt.export_service import export_project_to_pptx, get_export_file_path

logger = logging.getLogger(__name__)

ppt_export_bp = Blueprint("ppt_export", __name__)

# Export mode constants
EXPORT_MODE_SIMPLE = "simple"       # Images as backgrounds (default)
EXPORT_MODE_EDITABLE = "editable"   # Editable text boxes over clean background


def _get_upload_folder():
    """Get upload folder path"""
    return Path(current_app.config.get("UPLOAD_FOLDER", "uploads"))


@ppt_export_bp.route("/<project_id>/export", methods=["POST"])
@jwt_required()
def export_project(project_id: str):
    """
    POST /api/ppt/projects/{project_id}/export
    Export project to PPTX or PDF format.

    Request body:
        {
            "format": "pptx",           // pptx or pdf
            "mode": "simple",           // simple (images as backgrounds) or editable (text boxes)
            "filename": "optional_name"
        }

    Returns:
        {
            "success": true,
            "data": {
                "download_url": "/api/ppt/files/{project_id}/exports/{filename}",
                "filename": "presentation.pptx",
                "pages_exported": 10,
                "mode": "simple"
            }
        }
    """
    try:
        user_id = get_jwt_identity()

        # Verify project exists and user has access
        project = PPTProject.query.filter_by(id=project_id).first()
        if not project:
            return {"success": False, "error": "NOT_FOUND", "message": "Project not found"}, 404

        if project.user_id != user_id:
            return {"success": False, "error": "FORBIDDEN", "message": "Access denied"}, 403

        # Get request params
        data = request.get_json() or {}
        export_format = data.get("format", "pptx").lower()
        export_mode = data.get("mode", EXPORT_MODE_SIMPLE).lower()
        custom_filename = data.get("filename")

        if export_format not in ["pptx", "pdf"]:
            return {
                "success": False,
                "error": "INVALID_FORMAT",
                "message": "Format must be 'pptx' or 'pdf'"
            }, 400

        if export_mode not in [EXPORT_MODE_SIMPLE, EXPORT_MODE_EDITABLE]:
            return {
                "success": False,
                "error": "INVALID_MODE",
                "message": f"Mode must be '{EXPORT_MODE_SIMPLE}' or '{EXPORT_MODE_EDITABLE}'"
            }, 400

        # Get project pages ordered by index
        pages = PPTPage.query.filter_by(project_id=project_id).order_by(PPTPage.order_index).all()

        if not pages:
            return {
                "success": False,
                "error": "NO_PAGES",
                "message": "Project has no pages to export"
            }, 400

        # Check if pages have generated images
        pages_with_images = [p for p in pages if p.generated_image_path]
        if not pages_with_images:
            return {
                "success": False,
                "error": "NO_IMAGES",
                "message": "No generated images found. Please generate images first."
            }, 400

        upload_folder = _get_upload_folder()

        # Handle export based on mode
        if export_mode == EXPORT_MODE_EDITABLE:
            # Editable export - requires AI processing
            relative_path, export_info = _export_editable(
                project=project,
                pages=pages_with_images,
                upload_folder=upload_folder,
                user_id=user_id,
                filename=custom_filename,
                export_format=export_format
            )
        else:
            # Simple export - images as backgrounds
            relative_path = export_project_to_pptx(
                project_id=project_id,
                pages=pages_with_images,
                upload_folder=upload_folder,
                filename=custom_filename
            )
            export_info = {"text_elements": 0}

            # Convert to PDF if requested
            if export_format == "pdf":
                pptx_path = upload_folder / "ppt" / relative_path.lstrip('/')
                relative_path = _convert_pptx_to_pdf(pptx_path, upload_folder, project_id)

        # Extract filename from path
        filename = relative_path.split('/')[-1]

        return {
            "success": True,
            "data": {
                "download_url": f"/api/ppt/files{relative_path}",
                "filename": filename,
                "pages_exported": len(pages_with_images),
                "mode": export_mode,
                "format": export_format,
                **export_info
            }
        }

    except ValueError as e:
        logger.error(f"Export validation error: {str(e)}")
        return {"success": False, "error": "VALIDATION_ERROR", "message": str(e)}, 400
    except Exception as e:
        logger.error(f"Export failed: {str(e)}", exc_info=True)
        return {"success": False, "error": "SERVER_ERROR", "message": str(e)}, 500


def _export_editable(project, pages, upload_folder, user_id, filename, export_format):
    """
    Export using editable mode with AI-powered text extraction.

    Returns:
        Tuple of (relative_path, export_info_dict)
    """
    from backend.services.ppt.editable_export_service import create_editable_pptx
    from backend.services.ppt.file_service import get_ai_provider_for_user

    # Get AI provider
    ai_provider = get_ai_provider_for_user(user_id)
    if not ai_provider:
        raise ValueError("AI provider not configured. Please configure your AI settings.")

    # Prepare page images
    page_images = []
    page_descriptions = []

    for page in pages:
        image_path = page.generated_image_path
        if image_path.startswith('/'):
            image_path = image_path.lstrip('/')

        full_path = upload_folder / image_path
        if not full_path.exists():
            full_path = upload_folder / "ppt" / project.id / "pages" / image_path.split('/')[-1]

        if full_path.exists():
            with open(full_path, 'rb') as f:
                image_data = f.read()
            page_images.append((page.id, image_data))

            # Get description
            desc_content = page.get_description_content() if hasattr(page, 'get_description_content') else None
            desc_text = desc_content.get('text', '') if desc_content else ''
            page_descriptions.append(desc_text)

    if not page_images:
        raise ValueError("No images found to export")

    # Generate output path
    exports_dir = upload_folder / "ppt" / project.id / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)

    if not filename:
        filename = f"editable_{project.id[:8]}.pptx"
    elif not filename.endswith('.pptx'):
        filename += '.pptx'

    output_path = str(exports_dir / filename)

    # Get export settings from project
    inpaint_method = project.export_inpaint_method or "generative"
    extractor_method = project.export_extractor_method or "hybrid"

    # Create editable PPTX
    result = create_editable_pptx(
        ai_provider=ai_provider,
        project_id=project.id,
        page_images=page_images,
        output_path=output_path,
        inpaint_method=inpaint_method,
        extractor_method=extractor_method,
        page_descriptions=page_descriptions
    )

    if not result.success:
        raise ValueError(f"Editable export failed: {result.error}")

    relative_path = f"/{project.id}/exports/{filename}"

    # Convert to PDF if requested
    if export_format == "pdf":
        relative_path = _convert_pptx_to_pdf(Path(output_path), upload_folder, project.id)

    return relative_path, {"text_elements": result.text_elements_found}


def _convert_pptx_to_pdf(pptx_path: Path, upload_folder: Path, project_id: str) -> str:
    """
    Convert PPTX to PDF using LibreOffice.

    Args:
        pptx_path: Path to the PPTX file
        upload_folder: Base upload folder
        project_id: Project ID

    Returns:
        Relative path to the PDF file
    """
    if not pptx_path.exists():
        raise ValueError(f"PPTX file not found: {pptx_path}")

    exports_dir = upload_folder / "ppt" / project_id / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)

    pdf_filename = pptx_path.stem + ".pdf"
    pdf_path = exports_dir / pdf_filename

    try:
        # Try using LibreOffice for conversion
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(exports_dir),
                str(pptx_path)
            ],
            capture_output=True,
            text=True,
            timeout=120  # 2 minute timeout
        )

        if result.returncode != 0:
            logger.warning(f"LibreOffice conversion warning: {result.stderr}")

        if pdf_path.exists():
            return f"/{project_id}/exports/{pdf_filename}"
        else:
            raise ValueError("PDF conversion completed but file not found")

    except FileNotFoundError:
        # LibreOffice not installed, try alternative method
        logger.warning("LibreOffice not found, attempting alternative PDF conversion")

        try:
            # Try using pdf2pptx or similar Python libraries
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.pdfgen import canvas

            # Create PDF from PPTX slides
            prs = Presentation(str(pptx_path))

            # Use reportlab to create PDF
            pdf_canvas = canvas.Canvas(str(pdf_path), pagesize=landscape(A4))
            width, height = landscape(A4)

            for slide in prs.slides:
                # For each slide, extract images and add to PDF
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        # This is a simplified approach - full implementation would
                        # properly render all slide content
                        img_stream = shape.image.blob
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                            tmp.write(img_stream)
                            tmp_path = tmp.name

                        try:
                            pdf_canvas.drawImage(tmp_path, 0, 0, width=width, height=height)
                        finally:
                            Path(tmp_path).unlink(missing_ok=True)

                        break  # Only use first image per slide

                pdf_canvas.showPage()

            pdf_canvas.save()
            return f"/{project_id}/exports/{pdf_filename}"

        except ImportError:
            raise ValueError(
                "PDF conversion requires LibreOffice or reportlab. "
                "Please install LibreOffice for best results."
            )

    except subprocess.TimeoutExpired:
        raise ValueError("PDF conversion timed out")
    except Exception as e:
        raise ValueError(f"PDF conversion failed: {str(e)}")


@ppt_export_bp.route("/<project_id>/export/download/<filename>")
@jwt_required()
def download_export(project_id: str, filename: str):
    """
    GET /api/ppt/projects/{project_id}/export/download/{filename}
    Download an exported file (with auth).
    """
    try:
        user_id = get_jwt_identity()

        # Verify project access
        project = PPTProject.query.filter_by(id=project_id).first()
        if not project:
            return {"success": False, "error": "NOT_FOUND", "message": "Project not found"}, 404

        if project.user_id != user_id:
            return {"success": False, "error": "FORBIDDEN", "message": "Access denied"}, 403

        upload_folder = _get_upload_folder()
        filepath = get_export_file_path(project_id, filename, upload_folder)

        if not filepath:
            return {"success": False, "error": "NOT_FOUND", "message": "Export file not found"}, 404

        # Determine mimetype
        if filename.endswith('.pptx'):
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        elif filename.endswith('.pdf'):
            mimetype = 'application/pdf'
        else:
            mimetype = 'application/octet-stream'

        return send_file(
            filepath,
            mimetype=mimetype,
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        logger.error(f"Download failed: {str(e)}", exc_info=True)
        return {"success": False, "error": "SERVER_ERROR", "message": str(e)}, 500
