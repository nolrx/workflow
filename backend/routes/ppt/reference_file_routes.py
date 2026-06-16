"""
PPT Reference File routes - handles file upload and management
"""
import logging
import os
import re
import threading
import uuid
from datetime import datetime
from pathlib import Path
from urllib.parse import unquote

from flask import Blueprint, current_app, request, send_file
from flask_jwt_extended import get_jwt_identity, jwt_required
from werkzeug.utils import secure_filename

from backend.extensions import db
from backend.models.ppt import PPTProject, PPTReferenceFile
from backend.utils.response import error_response, success_response

logger = logging.getLogger(__name__)

ppt_reference_file_bp = Blueprint("ppt_reference_files", __name__)

# Allowed file extensions for reference files
ALLOWED_REFERENCE_EXTENSIONS = {"pdf", "docx", "pptx", "xlsx", "csv", "txt", "md"}


def _get_upload_folder():
    """Get upload folder path"""
    return Path(current_app.config.get("UPLOAD_FOLDER", "uploads"))


def _allowed_file(filename: str) -> bool:
    """Check if file extension is allowed"""
    return "." in filename and \
           filename.rsplit(".", 1)[1].lower() in ALLOWED_REFERENCE_EXTENSIONS


def _get_file_type(filename: str) -> str:
    """Get file type from filename"""
    if "." in filename:
        return filename.rsplit(".", 1)[1].lower()
    return "unknown"


def _parse_file_content(file_path: str, file_type: str) -> tuple[str, str]:
    """
    Parse file content to markdown (simplified version)

    Args:
        file_path: Path to the file
        file_type: File extension type

    Returns:
        (markdown_content, error_message)
    """
    try:
        if file_type in ["txt", "md"]:
            # Plain text files
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            return content, None

        elif file_type == "csv":
            # CSV files - convert to markdown table
            import csv
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                rows = list(reader)

            if not rows:
                return "", None

            # Create markdown table
            md_lines = []
            # Header
            md_lines.append("| " + " | ".join(rows[0]) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
            # Data rows
            for row in rows[1:]:
                md_lines.append("| " + " | ".join(row) + " |")

            return "\n".join(md_lines), None

        elif file_type == "docx":
            # Word documents - try to extract text
            try:
                from docx import Document
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return "\n\n".join(paragraphs), None
            except ImportError:
                return None, "python-docx library not installed for DOCX parsing"
            except Exception as e:
                return None, f"Failed to parse DOCX: {str(e)}"

        elif file_type == "pdf":
            # PDF files - note: full parsing requires external service
            return None, "PDF parsing requires external service configuration"

        elif file_type == "pptx":
            # PowerPoint files
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                content_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        content_parts.append(f"## Slide {slide_num}\n\n" + "\n\n".join(slide_text))
                return "\n\n---\n\n".join(content_parts), None
            except ImportError:
                return None, "python-pptx library not installed for PPTX parsing"
            except Exception as e:
                return None, f"Failed to parse PPTX: {str(e)}"

        elif file_type == "xlsx":
            # Excel files
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True)
                content_parts = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    rows = list(sheet.iter_rows(values_only=True))
                    if not rows:
                        continue

                    md_lines = [f"## {sheet_name}\n"]
                    # Find first row with data for header
                    header = rows[0] if rows else []
                    if header:
                        header_strs = [str(c) if c else "" for c in header]
                        md_lines.append("| " + " | ".join(header_strs) + " |")
                        md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

                        for row in rows[1:]:
                            row_strs = [str(c) if c else "" for c in row]
                            md_lines.append("| " + " | ".join(row_strs) + " |")

                    content_parts.append("\n".join(md_lines))
                wb.close()
                return "\n\n".join(content_parts), None
            except ImportError:
                return None, "openpyxl library not installed for XLSX parsing"
            except Exception as e:
                return None, f"Failed to parse XLSX: {str(e)}"

        else:
            return None, f"Unsupported file type: {file_type}"

    except Exception as e:
        return None, f"Error parsing file: {str(e)}"


def _parse_file_async(file_id: str, file_path: str, file_type: str, app):
    """
    Parse file asynchronously in background

    Args:
        file_id: Reference file ID
        file_path: Path to the uploaded file
        file_type: File type extension
        app: Flask app instance (for app context)
    """
    with app.app_context():
        try:
            reference_file = PPTReferenceFile.query.get(file_id)
            if not reference_file:
                logger.error(f"Reference file {file_id} not found")
                return

            # Update status to parsing
            reference_file.parse_status = "parsing"
            db.session.commit()

            # Parse file
            logger.info(f"Starting to parse file: {file_path}")
            markdown_content, error_message = _parse_file_content(file_path, file_type)

            # Update database
            if error_message:
                reference_file.parse_status = "failed"
                reference_file.error_message = error_message
                logger.error(f"File parsing failed: {error_message}")
            else:
                reference_file.parse_status = "completed"
                reference_file.markdown_content = markdown_content
                logger.info(f"File parsing completed: {file_path}")

            reference_file.updated_at = datetime.utcnow()
            db.session.commit()

        except Exception as e:
            logger.error(f"Error in async file parsing: {str(e)}", exc_info=True)
            try:
                reference_file = PPTReferenceFile.query.get(file_id)
                if reference_file:
                    reference_file.parse_status = "failed"
                    reference_file.error_message = f"Parsing error: {str(e)}"
                    reference_file.updated_at = datetime.utcnow()
                    db.session.commit()
            except Exception as db_error:
                logger.error(f"Failed to update error status: {str(db_error)}")


@ppt_reference_file_bp.route("/reference-files/upload", methods=["POST"])
@jwt_required()
def upload_reference_file():
    """
    POST /api/ppt/reference-files/upload - Upload a reference file

    Supports multipart/form-data:
    - file: The file to upload (required)
    - project_id: Project ID to associate with (optional)

    Returns:
        Reference file information with status
    """
    try:
        user_id = get_jwt_identity()

        # Check if file is in request
        if "file" not in request.files:
            return error_response("BAD_REQUEST", "No file provided", 400)

        file = request.files["file"]

        # Get filename - handle encoding issues
        original_filename = file.filename
        if not original_filename or original_filename == "":
            content_disposition = request.headers.get("Content-Disposition", "")
            if content_disposition:
                filename_match = re.search(r'filename[^;=\n]*=(([\'"]).*?\2|[^;\n]*)', content_disposition)
                if filename_match:
                    original_filename = filename_match.group(1).strip("\"'")
                    try:
                        original_filename = unquote(original_filename)
                    except Exception:
                        pass

        if not original_filename or original_filename == "":
            return error_response("BAD_REQUEST", "No file selected", 400)

        logger.info(f"Received file upload: {original_filename}")

        # Check file extension
        if not _allowed_file(original_filename):
            return error_response(
                "BAD_REQUEST",
                f"File type not allowed. Allowed types: {', '.join(ALLOWED_REFERENCE_EXTENSIONS)}",
                400
            )

        # Get project_id (optional)
        project_id = request.form.get("project_id")
        if project_id in ["none", ""]:
            project_id = None
        elif project_id:
            # Verify project exists and user has access
            project = PPTProject.query.get(project_id)
            if not project:
                return error_response("NOT_FOUND", "Project not found", 404)
            if project.user_id != user_id:
                return error_response("FORBIDDEN", "Access denied", 403)

        # Secure filename for filesystem
        filename = secure_filename(original_filename)

        # If secure_filename removed everything, use a fallback
        if not filename:
            ext = _get_file_type(original_filename)
            if ext == "unknown":
                ext = "file"
            filename = f"file_{uuid.uuid4().hex[:8]}.{ext}"
            logger.warning(f"Original filename '{original_filename}' was sanitized to '{filename}'")

        # Create upload directory
        upload_folder = _get_upload_folder()
        reference_files_dir = upload_folder / "ppt" / "reference_files"
        reference_files_dir.mkdir(parents=True, exist_ok=True)

        # Generate unique filename
        unique_id = str(uuid.uuid4())[:8]
        file_type = _get_file_type(original_filename)
        unique_filename = f"{unique_id}_{filename}"
        file_path = reference_files_dir / unique_filename

        # Save file
        file.save(str(file_path))
        file_size = os.path.getsize(file_path)

        # Create database record
        reference_file = PPTReferenceFile(
            user_id=user_id,
            project_id=project_id,
            filename=original_filename,
            file_path=str(file_path.relative_to(upload_folder)),
            file_size=file_size,
            file_type=file_type,
            parse_status="pending"
        )

        db.session.add(reference_file)
        db.session.commit()

        logger.info(f"File uploaded: {original_filename} (ID: {reference_file.id})")

        return success_response({"file": reference_file.to_dict(include_content=False)})

    except Exception as e:
        db.session.rollback()
        logger.error(f"Error uploading reference file: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/<file_id>", methods=["GET"])
@jwt_required()
def get_reference_file(file_id: str):
    """
    GET /api/ppt/reference-files/{file_id} - Get reference file information

    Returns:
        Reference file information including parse status and content
    """
    try:
        user_id = get_jwt_identity()

        reference_file = PPTReferenceFile.query.get(file_id)
        if not reference_file:
            return error_response("NOT_FOUND", "Reference file not found", 404)

        # Verify ownership
        if reference_file.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        return success_response({
            "file": reference_file.to_dict(include_content=True, include_failed_count=True)
        })

    except Exception as e:
        logger.error(f"Error getting reference file: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/<file_id>", methods=["DELETE"])
@jwt_required()
def delete_reference_file(file_id: str):
    """
    DELETE /api/ppt/reference-files/{file_id} - Delete a reference file

    Returns:
        Success message
    """
    try:
        user_id = get_jwt_identity()

        reference_file = PPTReferenceFile.query.get(file_id)
        if not reference_file:
            return error_response("NOT_FOUND", "Reference file not found", 404)

        # Verify ownership
        if reference_file.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        # Delete file from disk
        try:
            upload_folder = _get_upload_folder()
            file_path = upload_folder / reference_file.file_path
            if file_path.exists():
                file_path.unlink()
                logger.info(f"Deleted file from disk: {file_path}")
        except Exception as e:
            logger.warning(f"Failed to delete file from disk: {str(e)}")

        # Delete from database
        db.session.delete(reference_file)
        db.session.commit()

        logger.info(f"Deleted reference file: {file_id}")

        return success_response(message="File deleted successfully")

    except Exception as e:
        db.session.rollback()
        logger.error(f"Error deleting reference file: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/project/<project_id>", methods=["GET"])
@jwt_required()
def list_project_reference_files(project_id: str):
    """
    GET /api/ppt/reference-files/project/{project_id} - List reference files for a project

    Special values:
    - 'all': List all user's reference files
    - 'global': List only global files (not associated with any project)
    - project_id: List files for specific project

    Returns:
        List of reference files
    """
    try:
        user_id = get_jwt_identity()

        query = PPTReferenceFile.query.filter_by(user_id=user_id)

        if project_id == "all":
            # All user's files
            pass
        elif project_id in ["global", "none"]:
            # Global files only
            query = query.filter(PPTReferenceFile.project_id.is_(None))
        else:
            # Verify project exists and user has access
            project = PPTProject.query.get(project_id)
            if not project:
                return error_response("NOT_FOUND", "Project not found", 404)
            if project.user_id != user_id:
                # Check team membership
                if project.team_id:
                    from backend.models.team import TeamMember
                    membership = TeamMember.query.filter_by(
                        team_id=project.team_id, user_id=user_id
                    ).first()
                    if not membership:
                        return error_response("FORBIDDEN", "Access denied", 403)
                else:
                    return error_response("FORBIDDEN", "Access denied", 403)

            query = query.filter(PPTReferenceFile.project_id == project_id)

        reference_files = query.order_by(PPTReferenceFile.created_at.desc()).all()

        return success_response({
            "files": [f.to_dict(include_content=False) for f in reference_files]
        })

    except Exception as e:
        logger.error(f"Error listing reference files: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/<file_id>/parse", methods=["POST"])
@jwt_required()
def trigger_file_parse(file_id: str):
    """
    POST /api/ppt/reference-files/{file_id}/parse - Trigger parsing for a reference file

    Returns:
        Updated reference file information
    """
    try:
        user_id = get_jwt_identity()

        reference_file = PPTReferenceFile.query.get(file_id)
        if not reference_file:
            return error_response("NOT_FOUND", "Reference file not found", 404)

        # Verify ownership
        if reference_file.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        # If already parsing, return current status
        if reference_file.parse_status == "parsing":
            return success_response({
                "file": reference_file.to_dict(include_content=False),
                "message": "File is already being parsed"
            })

        # Reset for re-parsing if needed
        if reference_file.parse_status in ["completed", "failed"]:
            reference_file.parse_status = "pending"
            reference_file.error_message = None
            reference_file.markdown_content = None
            db.session.commit()

        # Get file path
        upload_folder = _get_upload_folder()
        file_path = upload_folder / reference_file.file_path

        if not file_path.exists():
            return error_response("FILE_NOT_FOUND", f"File not found: {file_path}", 404)

        # Start async parsing
        thread = threading.Thread(
            target=_parse_file_async,
            args=(
                reference_file.id,
                str(file_path),
                reference_file.file_type,
                current_app._get_current_object()
            )
        )
        thread.daemon = True
        thread.start()

        logger.info(f"Triggered parsing for file: {reference_file.filename} (ID: {file_id})")

        return success_response({
            "file": reference_file.to_dict(include_content=False),
            "message": "Parsing started"
        })

    except Exception as e:
        logger.error(f"Error triggering file parse: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/<file_id>/associate", methods=["POST"])
@jwt_required()
def associate_file_to_project(file_id: str):
    """
    POST /api/ppt/reference-files/{file_id}/associate - Associate a reference file to a project

    Request body:
    {
        "project_id": "project-id-here"
    }

    Returns:
        Updated reference file information
    """
    try:
        user_id = get_jwt_identity()

        reference_file = PPTReferenceFile.query.get(file_id)
        if not reference_file:
            return error_response("NOT_FOUND", "Reference file not found", 404)

        # Verify ownership
        if reference_file.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        data = request.get_json() or {}
        project_id = data.get("project_id")

        if not project_id:
            return error_response("BAD_REQUEST", "project_id is required", 400)

        # Verify project exists and user has access
        project = PPTProject.query.get(project_id)
        if not project:
            return error_response("NOT_FOUND", "Project not found", 404)
        if project.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        # Update file's project_id
        reference_file.project_id = project_id
        reference_file.updated_at = datetime.utcnow()
        db.session.commit()

        logger.info(f"Associated reference file {file_id} to project {project_id}")

        return success_response({"file": reference_file.to_dict(include_content=False)})

    except Exception as e:
        db.session.rollback()
        logger.error(f"Error associating reference file: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/<file_id>/dissociate", methods=["POST"])
@jwt_required()
def dissociate_file_from_project(file_id: str):
    """
    POST /api/ppt/reference-files/{file_id}/dissociate - Remove a reference file from its project

    Returns:
        Updated reference file information
    """
    try:
        user_id = get_jwt_identity()

        reference_file = PPTReferenceFile.query.get(file_id)
        if not reference_file:
            return error_response("NOT_FOUND", "Reference file not found", 404)

        # Verify ownership
        if reference_file.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        # Remove project association
        reference_file.project_id = None
        reference_file.updated_at = datetime.utcnow()
        db.session.commit()

        logger.info(f"Dissociated reference file {file_id} from project")

        return success_response({
            "file": reference_file.to_dict(include_content=False),
            "message": "File removed from project"
        })

    except Exception as e:
        db.session.rollback()
        logger.error(f"Error dissociating reference file: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)


@ppt_reference_file_bp.route("/reference-files/<file_id>/download", methods=["GET"])
@jwt_required()
def download_reference_file(file_id: str):
    """
    GET /api/ppt/reference-files/{file_id}/download - Download a reference file

    Returns:
        The file as attachment
    """
    try:
        user_id = get_jwt_identity()

        reference_file = PPTReferenceFile.query.get(file_id)
        if not reference_file:
            return error_response("NOT_FOUND", "Reference file not found", 404)

        # Verify ownership
        if reference_file.user_id != user_id:
            return error_response("FORBIDDEN", "Access denied", 403)

        upload_folder = _get_upload_folder()
        file_path = upload_folder / reference_file.file_path

        if not file_path.exists():
            return error_response("FILE_NOT_FOUND", "File not found on disk", 404)

        return send_file(
            str(file_path),
            as_attachment=True,
            download_name=reference_file.filename
        )

    except Exception as e:
        logger.error(f"Error downloading reference file: {str(e)}", exc_info=True)
        return error_response("SERVER_ERROR", str(e), 500)
